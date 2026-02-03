"""
Flash Reports Backend API

FastAPI application for generating HTML reports from PPTX templates.

Pipeline:
1. PPTX Template → PDF → PNG (conversion)
2. PNG → HTML Template (Claude Vision)
3. HTML Template + Project Data + Mapping → Populated HTML (Claude)
"""

from fastapi import FastAPI, HTTPException, BackgroundTasks, Header
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import HTMLResponse, JSONResponse
from pydantic import BaseModel
from typing import Optional, Dict, Any, List
import traceback
import time
from datetime import datetime

from app.config import SUPABASE_URL
from app.services.converter import convert_pptx_to_images
from app.services.claude_html import generate_html_template, extract_template_fields
from app.services.data_populator import (
    generate_multi_project_html,
    populate_html_with_claude,
    simple_populate_html,
    apply_mapping_to_project
)
# PDF generation - optional, requires system dependencies (GLib, Pango, Cairo)
PDF_GENERATION_AVAILABLE = False
html_to_pdf = None
try:
    from app.services.pdf_generator import html_to_pdf
    PDF_GENERATION_AVAILABLE = True
except (ImportError, OSError) as e:
    print(f"Warning: PDF generation not available (missing system libraries): {e}")
from app.services.supabase_client import (
    get_session,
    get_mapping,
    get_fetched_projects,
    download_template,
    upload_generated_html,
    upload_pdf,
    upload_pptx,
    upload_png,
    create_generation_job,
    update_job_status,
    save_generated_report,
    get_job_status,
    update_template_preparation_status,
    get_template_preparation_status,
    download_html_template
)

# PPTX generation from HTML
PPTX_GENERATION_AVAILABLE = False
html_to_pptx_convert = None
try:
    from app.services.pptx_generator import html_to_pptx
    html_to_pptx_convert = html_to_pptx
    PPTX_GENERATION_AVAILABLE = True
except (ImportError, OSError) as e:
    print(f"Warning: PPTX generation from HTML not available: {e}")


app = FastAPI(
    title="Flash Reports API",
    description="Generate HTML reports from PPTX templates using Claude Vision",
    version="1.0.0"
)

# CORS configuration
app.add_middleware(
    CORSMiddleware,
    allow_origins=[
        "https://flash-reports-poc.vercel.app",
        "http://localhost:5173",
        "http://localhost:3000",
    ],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# Request/Response models
class GenerateJobRequest(BaseModel):
    use_claude_population: bool = True


class GenerateJobResponse(BaseModel):
    success: bool
    jobId: Optional[str] = None
    error: Optional[str] = None


class JobStatusRequest(BaseModel):
    job_id: str


class JobStatusResponse(BaseModel):
    success: bool
    job: Optional[Dict[str, Any]] = None
    error: Optional[str] = None


class TemplateAnalysisResponse(BaseModel):
    success: bool
    html_template: Optional[str] = None
    fields: Optional[List[Dict[str, Any]]] = None
    slide_count: Optional[int] = None
    error: Optional[str] = None


@app.get("/")
async def root():
    """API info endpoint."""
    return {
        "name": "Flash Reports API",
        "version": "1.0.0",
        "description": "Generate HTML reports from PPTX templates using Claude Vision",
        "pipeline": "PPTX → PDF → PNG → HTML Template → Populated HTML"
    }


@app.get("/health")
async def health_check():
    """Health check endpoint."""
    return {"status": "healthy", "timestamp": datetime.utcnow().isoformat()}


@app.post("/list-slides")
async def list_slides(
    x_session_id: str = Header(..., alias="x-session-id")
):
    """
    List all slides in the uploaded PPTX template with titles.
    Uses python-pptx for fast local extraction (no AI needed).
    """
    try:
        from pptx import Presentation
        import io

        session = await get_session(x_session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")

        template_path = session.get('template_path')
        if not template_path:
            try:
                mapping = await get_mapping(x_session_id)
                if mapping:
                    template_path = mapping.get('template_path')
            except Exception:
                pass
        if not template_path:
            raise HTTPException(status_code=400, detail="No template uploaded for this session")

        # Download PPTX from storage
        pptx_bytes = await download_template(template_path)
        prs = Presentation(io.BytesIO(pptx_bytes))

        slides = []
        for i, slide in enumerate(prs.slides):
            title = ""
            # Try to get title from title placeholder
            if slide.shapes.title and slide.shapes.title.text:
                title = slide.shapes.title.text.strip()
            else:
                # Fallback: get first text shape
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        title = shape.text_frame.text.strip()[:80]
                        break

            layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
            shape_count = len(slide.shapes)

            slides.append({
                "slide_number": i + 1,
                "title": title or f"Slide {i + 1}",
                "layout": layout_name,
                "shape_count": shape_count,
            })

        return JSONResponse(content={
            "success": True,
            "slides": slides,
            "total": len(slides),
        })
    except HTTPException:
        raise
    except Exception as e:
        print(f"Error listing slides: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/generate-html", response_model=GenerateJobResponse)
async def generate_html_report(
    request: GenerateJobRequest,
    background_tasks: BackgroundTasks,
    x_session_id: str = Header(..., alias="x-session-id")
):
    """
    Create a generation job and start processing in the background.

    This endpoint returns immediately with a job ID. Use /job-status to poll for completion.
    """
    try:
        # Create the job
        job_id = await create_generation_job(
            session_id=x_session_id,
            engine="claude-html",
            input_data={"use_claude_population": request.use_claude_population}
        )

        # Start background processing
        background_tasks.add_task(
            process_html_generation,
            job_id=job_id,
            session_id=x_session_id,
            use_claude_population=request.use_claude_population
        )

        return GenerateJobResponse(success=True, jobId=job_id)

    except Exception as e:
        traceback.print_exc()
        return GenerateJobResponse(success=False, error=str(e))


@app.post("/job-status", response_model=JobStatusResponse)
async def check_job_status_endpoint(
    request: JobStatusRequest,
    x_session_id: str = Header(..., alias="x-session-id")
):
    """
    Check the status of a generation job.
    """
    try:
        job = await get_job_status(request.job_id)

        if not job:
            return JobStatusResponse(success=False, error="Job not found")

        return JobStatusResponse(success=True, job=job)

    except Exception as e:
        return JobStatusResponse(success=False, error=str(e))


@app.post("/analyze-template", response_model=TemplateAnalysisResponse)
async def analyze_template(
    x_session_id: str = Header(..., alias="x-session-id")
):
    """
    Analyze a PPTX template and generate HTML template with placeholders.

    This converts the PPTX to images and uses Claude Vision to create an HTML template.
    """
    try:
        # Get session to find template path
        session = await get_session(x_session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")

        template_path = session.get('template_path')
        if not template_path:
            raise HTTPException(status_code=400, detail="No template uploaded for this session")

        print(f"[analyze-template] Downloading template: {template_path}")

        # Download template from Supabase Storage
        pptx_bytes = await download_template(template_path)

        print(f"[analyze-template] Converting PPTX to images...")

        # Convert PPTX to images
        images = convert_pptx_to_images(pptx_bytes)
        print(f"[analyze-template] Generated {len(images)} slide images")

        print(f"[analyze-template] Generating HTML template with Claude Vision...")

        # Generate HTML template using Claude Vision
        result = generate_html_template(images)

        html_template = result["full_html"]
        fields = result.get("fields", [])

        # If no fields returned by Claude, extract them from HTML
        if not fields:
            fields = extract_template_fields(html_template)

        print(f"[analyze-template] Generated HTML with {len(fields)} template fields")

        return TemplateAnalysisResponse(
            success=True,
            html_template=html_template,
            fields=fields,
            slide_count=len(images)
        )

    except Exception as e:
        traceback.print_exc()
        return TemplateAnalysisResponse(success=False, error=str(e))


@app.get("/preview-template", response_class=HTMLResponse)
async def preview_template(
    x_session_id: str = Header(..., alias="x-session-id")
):
    """
    Preview the HTML template (before data population).
    """
    try:
        session = await get_session(x_session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")

        # Check if we have a cached HTML template
        html_template = session.get('html_template')

        if not html_template:
            raise HTTPException(
                status_code=400,
                detail="No HTML template generated. Call /analyze-template first."
            )

        return HTMLResponse(content=html_template)

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


async def process_html_generation(
    job_id: str,
    session_id: str,
    use_claude_population: bool = True
):
    """
    Background task to process HTML generation.

    Steps:
    1. Get session data (template path, mapping, fetched projects)
    2. Download and convert PPTX template to images
    3. Generate HTML template using Claude Vision
    4. Populate HTML with project data using mapping
    5. Save result and update job status
    """
    job_start_time = time.time()
    step_times = {}

    try:
        # STEP 1: Load session data
        step_start = time.time()
        print(f"[STEP 1/6] Starting HTML generation for job {job_id}")
        await update_job_status(job_id, "processing")

        # Get session data
        session = await get_session(session_id)
        if not session:
            raise RuntimeError("Session not found")

        # Get mapping
        mapping = await get_mapping(session_id)
        if not mapping:
            raise RuntimeError("No mapping found for this session")

        mapping_json = mapping.get('mapping_json', {})
        long_text_strategy = mapping.get('long_text_strategy', 'summarize')

        # Get template path - check session first, then mapping (for reused mappings)
        template_path = session.get('template_path') or mapping.get('template_path')
        if not template_path:
            raise RuntimeError("No template uploaded for this session")

        # Get fetched projects data
        fetched_data = await get_fetched_projects(session_id)
        if not fetched_data:
            raise RuntimeError("No project data found for this session")

        projects = fetched_data.get('projects', [])
        if not projects:
            raise RuntimeError("No projects in fetched data")

        step_times['step1'] = time.time() - step_start
        print(f"         Completed in {step_times['step1']:.2f}s")

        # STEP 2: Validate projects
        step_start = time.time()
        print(f"[STEP 2/6] Found {len(projects)} projects to process")
        step_times['step2'] = time.time() - step_start
        print(f"         Completed in {step_times['step2']:.2f}s")

        # Check if we have pre-generated HTML template (from background preparation)
        html_template_url = session.get('html_template_url')
        template_preparation_status = session.get('template_preparation_status')
        pdf_url = session.get('template_pdf_url')

        if html_template_url and template_preparation_status == 'completed':
            # USE PRE-GENERATED HTML TEMPLATE (optimized path)
            step_start = time.time()
            print(f"[STEP 3-5] Using pre-generated HTML template (skipping PPTX conversion)")
            print(f"         HTML URL: {html_template_url}")

            # Download HTML template
            html_template = await download_html_template(html_template_url)
            print(f"         Downloaded HTML template ({len(html_template)} chars)")

            # Count slides from HTML
            import re
            slide_matches = re.findall(r'<div[^>]*class="[^"]*slide[^"]*"[^>]*data-slide-number="(\d+)"', html_template)
            slide_count = len(slide_matches) if slide_matches else 1
            print(f"         Found {slide_count} slides in HTML template")

            template_html_url = html_template_url
            step_times['step3_5'] = time.time() - step_start
            print(f"         Completed in {step_times['step3_5']:.2f}s")

        else:
            # FALLBACK: Convert PPTX on-the-fly (backward compatibility)
            print(f"[STEP 3-5] No pre-generated HTML, converting PPTX on-the-fly...")

            # STEP 3: Download template
            step_start = time.time()
            print(f"[STEP 3/6] Downloading template: {template_path}")
            pptx_bytes = await download_template(template_path)
            step_times['step3'] = time.time() - step_start
            print(f"         Completed in {step_times['step3']:.2f}s")

            # STEP 4: Convert PPTX to images (and get PDF)
            step_start = time.time()
            print(f"[STEP 4/6] Converting PPTX to images...")
            images, pdf_bytes = convert_pptx_to_images(pptx_bytes, return_pdf=True)
            slide_count = len(images)
            print(f"         Generated {slide_count} slide images")

            # Upload PDF to storage
            timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
            if pdf_bytes:
                pdf_filename = f"template_{timestamp}.pdf"
                pdf_url = await upload_pdf(session_id, pdf_bytes, pdf_filename)
                print(f"         Uploaded PDF to: {pdf_url}")

            step_times['step4'] = time.time() - step_start
            print(f"         Completed in {step_times['step4']:.2f}s")

            # STEP 5: Generate HTML template
            step_start = time.time()
            print(f"[STEP 5/6] Generating HTML template with Claude Vision...")
            print(f"         Using {len(mapping_json)} fields from user mapping")
            print(f"         Long text strategy: {long_text_strategy}")
            template_result = generate_html_template(images, mapping_json, long_text_strategy)
            html_template = template_result["full_html"]
            print(f"         Generated HTML template with {len(template_result.get('fields', []))} fields")

            # Upload HTML template (with placeholders) to storage
            template_html_filename = f"template_{timestamp}.html"
            template_html_url = await upload_generated_html(session_id, html_template, template_html_filename)
            print(f"         Uploaded HTML template to: {template_html_url}")
            step_times['step5'] = time.time() - step_start
            print(f"         Completed in {step_times['step5']:.2f}s")

        # Ensure timestamp is defined for filenames below
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")

        # STEP 6: Populate with data
        step_start = time.time()
        print(f"[STEP 6/6] Populating HTML with project data...")
        if len(projects) > 1:
            # Multiple projects - generate slides for each
            final_html = generate_multi_project_html(
                html_template,
                projects,
                mapping_json,
                use_claude=use_claude_population,
                long_text_strategy=long_text_strategy
            )
        else:
            # Single project
            if use_claude_population:
                final_html = populate_html_with_claude(
                    html_template,
                    projects[0],
                    mapping_json,
                    long_text_strategy=long_text_strategy
                )
            else:
                import re
                template_fields = list(set(re.findall(r'\{\{(\w+)\}\}', html_template)))
                field_values = apply_mapping_to_project(projects[0], mapping_json, template_fields)
                final_html = simple_populate_html(html_template, field_values)

        print(f"         HTML population complete")

        # Save HTML to storage
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        filename = f"report_{timestamp}.html"

        html_url = await upload_generated_html(session_id, final_html, filename)
        print(f"         Uploaded HTML to: {html_url}")

        # Generate PDF from HTML (optional - depends on system libraries)
        report_pdf_url = None
        report_pdf_storage_path = None
        if PDF_GENERATION_AVAILABLE and html_to_pdf:
            print(f"         Converting HTML to PDF...")
            try:
                report_pdf_bytes = html_to_pdf(final_html)
                report_pdf_filename = f"report_{timestamp}.pdf"
                report_pdf_url = await upload_pdf(session_id, report_pdf_bytes, report_pdf_filename)
                report_pdf_storage_path = f"{session_id}/{report_pdf_filename}"
                print(f"         Uploaded PDF to: {report_pdf_url}")
            except Exception as pdf_error:
                print(f"         Warning: PDF generation failed: {pdf_error}")
                report_pdf_url = None
        else:
            print(f"         PDF generation skipped (dependencies not available)")

        # Generate PPTX from HTML (convert HTML slides to editable PowerPoint)
        report_pptx_url = None
        if PPTX_GENERATION_AVAILABLE and html_to_pptx_convert:
            print(f"         Converting HTML to PPTX...")
            try:
                report_pptx_bytes = html_to_pptx_convert(final_html)
                report_pptx_filename = f"report_{timestamp}.pptx"
                report_pptx_url = await upload_pptx(session_id, report_pptx_bytes, report_pptx_filename)
                print(f"         Uploaded PPTX to: {report_pptx_url}")
            except Exception as pptx_error:
                print(f"         Warning: PPTX generation failed: {pptx_error}")
                report_pptx_url = None
        else:
            print(f"         PPTX generation skipped (dependencies not available)")

        # Save report reference
        report_id = await save_generated_report(session_id, html_url, "claude-html")
        step_times['step6'] = time.time() - step_start
        print(f"         Completed in {step_times['step6']:.2f}s")

        # Update job as completed
        await update_job_status(
            job_id,
            "completed",
            result={
                "reportId": report_id,
                "htmlUrl": html_url,
                "pdfUrl": report_pdf_url,
                "pdfStoragePath": report_pdf_storage_path,
                "pptxUrl": report_pptx_url,
                "templateHtmlUrl": template_html_url,
                "templatePdfUrl": pdf_url,
                "projectCount": len(projects),
                "slideCount": slide_count
            }
        )

        total_time = time.time() - job_start_time
        print(f"\n{'='*50}")
        print(f"JOB TIMING SUMMARY for {job_id}")
        print(f"{'='*50}")
        print(f"  Step 1 (Load session data):     {step_times.get('step1', 0):.2f}s")
        print(f"  Step 2 (Validate projects):     {step_times.get('step2', 0):.2f}s")
        print(f"  Step 3 (Download template):     {step_times.get('step3', 0):.2f}s")
        print(f"  Step 4 (Convert PPTX to PNG):   {step_times.get('step4', 0):.2f}s")
        print(f"  Step 5 (Claude Vision HTML):    {step_times.get('step5', 0):.2f}s")
        print(f"  Step 6 (Claude Population):     {step_times.get('step6', 0):.2f}s")
        print(f"{'='*50}")
        print(f"  TOTAL TIME:                     {total_time:.2f}s ({total_time/60:.1f} min)")
        print(f"{'='*50}")
        print(f"✅ Job {job_id} completed successfully")

    except Exception as e:
        error_msg = str(e)
        total_time = time.time() - job_start_time
        print(f"❌ Job {job_id} failed after {total_time:.2f}s: {error_msg}")
        traceback.print_exc()

        await update_job_status(job_id, "failed", error=error_msg)


# Direct generation endpoint (for testing without job system)
@app.post("/generate-direct", response_class=HTMLResponse)
async def generate_direct(
    x_session_id: str = Header(..., alias="x-session-id"),
    use_claude: bool = True
):
    """
    Synchronous endpoint for direct HTML generation (testing only).

    Warning: This can take several minutes depending on template complexity.
    Use /generate-html for production.
    """
    try:
        # Get all required data
        session = await get_session(x_session_id)
        if not session:
            raise HTTPException(status_code=404, detail="Session not found")

        template_path = session.get('template_path')
        if not template_path:
            raise HTTPException(status_code=400, detail="No template uploaded")

        mapping = await get_mapping(x_session_id)
        if not mapping:
            raise HTTPException(status_code=400, detail="No mapping found")

        fetched_data = await get_fetched_projects(x_session_id)
        if not fetched_data or not fetched_data.get('projects'):
            raise HTTPException(status_code=400, detail="No project data found")

        # Download and convert
        pptx_bytes = await download_template(template_path)
        images = convert_pptx_to_images(pptx_bytes)

        # Get mapping for template generation
        mapping_json = mapping.get('mapping_json', {})
        long_text_strategy = mapping.get('long_text_strategy', 'summarize')

        # Generate template with user's field names and long text strategy
        template_result = generate_html_template(images, mapping_json, long_text_strategy)
        html_template = template_result["full_html"]

        # Populate
        projects = fetched_data['projects']

        if len(projects) > 1:
            final_html = generate_multi_project_html(
                html_template, projects, mapping_json, use_claude=use_claude,
                long_text_strategy=long_text_strategy
            )
        else:
            if use_claude:
                final_html = populate_html_with_claude(
                    html_template, projects[0], mapping_json,
                    long_text_strategy=long_text_strategy
                )
            else:
                import re
                template_fields = list(set(re.findall(r'\{\{(\w+)\}\}', html_template)))
                field_values = apply_mapping_to_project(projects[0], mapping_json, template_fields)
                final_html = simple_populate_html(html_template, field_values)

        return HTMLResponse(content=final_html)

    except HTTPException:
        raise
    except Exception as e:
        traceback.print_exc()
        raise HTTPException(status_code=500, detail=str(e))


# ==============================================================================
# TEMPLATE PREPARATION ENDPOINTS (Background PPTX → HTML conversion)
# ==============================================================================

class PrepareTemplateResponse(BaseModel):
    success: bool
    message: Optional[str] = None
    error: Optional[str] = None


class TemplatePreparationStatusResponse(BaseModel):
    success: bool
    status: str  # 'pending' | 'processing' | 'completed' | 'failed'
    htmlTemplateUrl: Optional[str] = None
    templatePngUrls: Optional[List[str]] = None
    templatePdfUrl: Optional[str] = None
    error: Optional[str] = None


class ListSlidesFromHtmlResponse(BaseModel):
    success: bool
    slides: Optional[List[Dict[str, Any]]] = None
    total: Optional[int] = None
    error: Optional[str] = None


@app.post("/prepare-template", response_model=PrepareTemplateResponse)
async def prepare_template(
    background_tasks: BackgroundTasks,
    x_session_id: str = Header(..., alias="x-session-id")
):
    """
    Start template preparation in background.

    Converts PPTX → PDF → PNG → HTML Template using Claude Vision.
    The HTML template is stored in Supabase Storage for later use.
    """
    try:
        # Get session to find template path
        session = await get_session(x_session_id)
        if not session:
            return PrepareTemplateResponse(success=False, error="Session not found")

        template_path = session.get('template_path')
        if not template_path:
            return PrepareTemplateResponse(success=False, error="No template uploaded for this session")

        # Check if already prepared for this template
        current_status = session.get('template_preparation_status')
        current_html_url = session.get('html_template_url')

        if current_status == 'completed' and current_html_url:
            # Already prepared
            return PrepareTemplateResponse(
                success=True,
                message="Template already prepared"
            )

        if current_status == 'processing':
            # Already in progress
            return PrepareTemplateResponse(
                success=True,
                message="Template preparation already in progress"
            )

        # Start background processing
        background_tasks.add_task(
            process_template_preparation,
            session_id=x_session_id,
            template_path=template_path
        )

        return PrepareTemplateResponse(
            success=True,
            message="Template preparation started"
        )

    except Exception as e:
        traceback.print_exc()
        return PrepareTemplateResponse(success=False, error=str(e))


async def process_template_preparation(session_id: str, template_path: str):
    """
    Background task to convert PPTX → PDF → PNG → HTML.

    Steps:
    1. Update status to 'processing'
    2. Download PPTX from Storage
    3. Convert PPTX → PDF → PNG
    4. Generate HTML template with Claude Vision
    5. Upload HTML, PNGs, PDF to Storage
    6. Update session with URLs and status
    """
    start_time = time.time()

    try:
        print(f"[prepare-template] Starting for session {session_id}")

        # Step 1: Mark as processing
        await update_template_preparation_status(session_id, 'processing')

        # Step 2: Download template
        print(f"[prepare-template] Downloading template: {template_path}")
        pptx_bytes = await download_template(template_path)

        # Step 3: Convert to images (and get PDF)
        print(f"[prepare-template] Converting PPTX to images...")
        images, pdf_bytes = convert_pptx_to_images(pptx_bytes, return_pdf=True)
        print(f"[prepare-template] Generated {len(images)} slide images")

        # Step 4: Generate HTML template with Claude Vision
        print(f"[prepare-template] Generating HTML template with Claude Vision...")
        template_result = generate_html_template(images)
        html_template = template_result["full_html"]
        print(f"[prepare-template] Generated HTML with {len(template_result.get('fields', []))} fields")

        # Step 5: Upload files to Storage
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")

        # Upload HTML
        html_filename = f"template_{timestamp}.html"
        html_url = await upload_generated_html(session_id, html_template, html_filename)
        print(f"[prepare-template] Uploaded HTML to: {html_url}")

        # Upload PNGs
        png_urls = []
        for i, (img_bytes, _) in enumerate(images):
            png_filename = f"slide_{timestamp}_{i+1:02d}.png"
            png_url = await upload_png(session_id, img_bytes, png_filename)
            png_urls.append(png_url)
        print(f"[prepare-template] Uploaded {len(png_urls)} PNG images")

        # Upload PDF
        pdf_url = None
        if pdf_bytes:
            pdf_filename = f"template_{timestamp}.pdf"
            pdf_url = await upload_pdf(session_id, pdf_bytes, pdf_filename)
            print(f"[prepare-template] Uploaded PDF to: {pdf_url}")

        # Step 6: Update session with URLs
        await update_template_preparation_status(
            session_id,
            'completed',
            html_template_url=html_url,
            template_png_urls=png_urls,
            template_pdf_url=pdf_url
        )

        elapsed = time.time() - start_time
        print(f"[prepare-template] ✅ Completed in {elapsed:.1f}s for session {session_id}")

    except Exception as e:
        elapsed = time.time() - start_time
        error_msg = str(e)
        print(f"[prepare-template] ❌ Failed after {elapsed:.1f}s: {error_msg}")
        traceback.print_exc()

        await update_template_preparation_status(
            session_id,
            'failed',
            error=error_msg
        )


@app.post("/template-preparation-status", response_model=TemplatePreparationStatusResponse)
async def check_template_preparation_status(
    x_session_id: str = Header(..., alias="x-session-id")
):
    """
    Check the status of template preparation.
    """
    try:
        status_info = await get_template_preparation_status(x_session_id)

        return TemplatePreparationStatusResponse(
            success=True,
            status=status_info.get('status', 'pending'),
            htmlTemplateUrl=status_info.get('html_template_url'),
            templatePngUrls=status_info.get('template_png_urls'),
            templatePdfUrl=status_info.get('template_pdf_url'),
            error=status_info.get('error')
        )

    except Exception as e:
        return TemplatePreparationStatusResponse(
            success=False,
            status='failed',
            error=str(e)
        )


@app.post("/list-slides-from-html", response_model=ListSlidesFromHtmlResponse)
async def list_slides_from_html(
    x_session_id: str = Header(..., alias="x-session-id")
):
    """
    List all slides from the prepared HTML template.

    Parses the HTML to extract slide information (number, title, field count).
    This is used by the SlideSelector component.
    """
    try:
        from bs4 import BeautifulSoup
        import re

        # Get session to check preparation status
        session = await get_session(x_session_id)
        if not session:
            return ListSlidesFromHtmlResponse(success=False, error="Session not found")

        # Check if template is prepared
        status = session.get('template_preparation_status')
        html_url = session.get('html_template_url')

        if status != 'completed' or not html_url:
            return ListSlidesFromHtmlResponse(
                success=False,
                error=f"Template not ready. Status: {status}"
            )

        # Download HTML
        html_content = await download_html_template(html_url)

        # Parse HTML
        soup = BeautifulSoup(html_content, 'html.parser')

        # Find all slide divs
        slides = []
        slide_divs = soup.find_all('div', class_='slide')

        for slide_div in slide_divs:
            # Get slide number from data attribute or position
            slide_number = slide_div.get('data-slide-number')
            if slide_number:
                slide_number = int(slide_number)
            else:
                slide_number = len(slides) + 1

            # Extract title (look for .main-title or first heading)
            title_elem = slide_div.find(class_='main-title')
            if not title_elem:
                title_elem = slide_div.find(['h1', 'h2', 'h3'])

            title = title_elem.get_text(strip=True) if title_elem else f"Slide {slide_number}"
            # Truncate long titles
            if len(title) > 60:
                title = title[:57] + "..."

            # Count content sections (section-box elements)
            section_boxes = slide_div.find_all(class_='section-box')
            content_count = len(section_boxes) if section_boxes else 1

            slides.append({
                "slide_number": slide_number,
                "title": title,
                "field_count": content_count,  # Number of content sections
                "layout": "content"  # Could be enhanced to detect layout type
            })

        return ListSlidesFromHtmlResponse(
            success=True,
            slides=slides,
            total=len(slides)
        )

    except Exception as e:
        traceback.print_exc()
        return ListSlidesFromHtmlResponse(success=False, error=str(e))


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
