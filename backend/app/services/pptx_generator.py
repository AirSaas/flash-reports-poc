"""
HTML-to-PPTX converter service.

Converts populated HTML slide presentations into faithful python-pptx files.
Fully generic: reads all colors, fonts, sizes, and positions from the HTML/CSS.
No hardcoded theme colors — works with any HTML color scheme.

Based on: https://github.com/damianstetsonair/html-to-pptx
"""

import re
import io
import tempfile
from typing import Optional, List, Tuple, Dict

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree, html as lxml_html

# ───────────────────────── constants ──────────────────────────
SLIDE_W_PX, SLIDE_H_PX = 960, 540
SLIDE_W_IN, SLIDE_H_IN = 10.0, 5.625
_SCALE = SLIDE_W_IN / SLIDE_W_PX  # inches per CSS pixel

# Fallback defaults (used only when CSS provides nothing)
_FALLBACK_FONT     = 'Arial'
_FALLBACK_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
_FALLBACK_BLACK33  = RGBColor(0x33, 0x33, 0x33)
_FALLBACK_GREY66   = RGBColor(0x66, 0x66, 0x66)
_FALLBACK_GREY_CC  = RGBColor(0xCC, 0xCC, 0xCC)
_FALLBACK_TEAL     = RGBColor(0x00, 0x62, 0x72)
_FALLBACK_RED_BULL = RGBColor(0xCC, 0x00, 0x00)

# ───────────────────────── helpers ────────────────────────────
def E(px: float) -> int:
    """CSS pixels → EMU."""
    return int(round(px * _SCALE * 914400))

def _parse_color(s: str) -> Optional[RGBColor]:
    if not s:
        return None
    s = s.strip().lower()
    # named
    _NAMED = {'white': _FALLBACK_WHITE, 'black': RGBColor(0,0,0),
              'red': RGBColor(0xFF,0,0), 'green': RGBColor(0,0x80,0),
              'blue': RGBColor(0,0,0xFF), 'transparent': None}
    if s in _NAMED:
        return _NAMED[s]
    # hex 6
    m = re.match(r'#([0-9a-f]{6})$', s)
    if m:
        h = m.group(1); return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
    # hex 3
    m = re.match(r'#([0-9a-f]{3})$', s)
    if m:
        h = m.group(1); return RGBColor(int(h[0]*2,16), int(h[1]*2,16), int(h[2]*2,16))
    # rgb()
    m = re.match(r'rgb\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)', s)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    # rgba() — ignore alpha
    m = re.match(r'rgba\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)', s)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    return None

def _is_light(c: RGBColor) -> bool:
    """True if colour has high luminance (light background → keep dark text)."""
    return (0.299 * c[0] + 0.587 * c[1] + 0.114 * c[2]) > 180

def _parse_border_color(border_str: str) -> Optional[RGBColor]:
    """Extract colour from a CSS border shorthand like '1px solid #ccc'."""
    if not border_str:
        return None
    # try last token
    parts = border_str.strip().split()
    for p in reversed(parts):
        c = _parse_color(p)
        if c: return c
    return None

def _sty(el) -> dict:
    """Parse inline style attribute → dict."""
    d = {}
    for part in (el.get('style', '') or '').split(';'):
        if ':' not in part: continue
        k, v = part.split(':', 1)
        d[k.strip().lower()] = v.strip()
    return d

def _px(v: str) -> float:
    m = re.match(r'([\d.]+)', v.strip()); return float(m.group(1)) if m else 0

def _pct(v: str) -> float:
    m = re.match(r'([\d.]+)\s*%', v.strip()); return float(m.group(1)) if m else 0

# ───────────── stylesheet parser ────────────────────────────
def _parse_stylesheet(doc) -> Dict[str, dict]:
    """Parse <style> blocks into {selector: {prop: value}} dict."""
    ss: Dict[str, dict] = {}
    for style_el in doc.cssselect('style'):
        raw = style_el.text_content() or ''
        raw = re.sub(r'/\*.*?\*/', '', raw, flags=re.DOTALL)
        for m in re.finditer(r'([^{}]+)\{([^}]*)\}', raw):
            selectors = m.group(1).strip()
            body = m.group(2).strip()
            props = {}
            for decl in body.split(';'):
                if ':' not in decl: continue
                k, v = decl.split(':', 1)
                props[k.strip().lower()] = v.strip()
            for sel in selectors.split(','):
                ss[sel.strip()] = props
    return ss

def _ss_get(ss: dict, *selectors) -> dict:
    """Merge stylesheet rules for the given selectors (later wins)."""
    merged = {}
    for sel in selectors:
        if sel in ss:
            merged.update(ss[sel])
    return merged

def _bg_color(sty: dict) -> Optional[RGBColor]:
    """Extract background color from a style dict."""
    return _parse_color(sty.get('background', '') or sty.get('background-color', ''))

def _resolve_width(sty: dict, key: str = 'width', default: float = 0, container_w: float = SLIDE_W_PX) -> float:
    """Resolve a CSS width value (px or %) to pixels."""
    val = sty.get(key, '')
    if not val:
        return default
    px_val = _px(val)
    if px_val > 0:
        return px_val
    if '%' in val:
        return container_w * _pct(val) / 100
    return default

# ───────────── XML-level table cell helpers ───────────────────
def _cell_border(cell, color=_FALLBACK_GREY_CC, width=6350, dash='solid'):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    for tag_name in ('a:lnT','a:lnB','a:lnL','a:lnR'):
        tag = qn(tag_name)
        for old in tcPr.findall(tag): tcPr.remove(old)
        ln = etree.SubElement(tcPr, tag)
        ln.set('w', str(int(width))); ln.set('cap','flat'); ln.set('cmpd','sng'); ln.set('algn','ctr')
        sf = etree.SubElement(ln, qn('a:solidFill'))
        srgb = etree.SubElement(sf, qn('a:srgbClr'))
        srgb.set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))
        if dash == 'dashed':
            etree.SubElement(ln, qn('a:prstDash')).set('val', 'dash')

def _cell_fill(cell, color: RGBColor):
    tc = cell._tc; tcPr = tc.get_or_add_tcPr()
    for old in tcPr.findall(qn('a:solidFill')): tcPr.remove(old)
    sf = etree.SubElement(tcPr, qn('a:solidFill'))
    etree.SubElement(sf, qn('a:srgbClr')).set('val', '%02X%02X%02X' % (color[0], color[1], color[2]))
    tcPr.insert(0, sf)

def _nuke_table_theme(shape):
    """Remove built-in table theme so manual cell fills/borders show."""
    tbl = shape._element.find('.//' + qn('a:tbl'))
    if tbl is None: return
    tblPr = tbl.find(qn('a:tblPr'))
    if tblPr is None: return
    for k in ('bandRow','bandCol','firstRow','lastRow','firstCol','lastCol'):
        tblPr.set(k, '0')
    for child_tag in (qn('a:tblStyle'), qn('a:tableStyleId')):
        for ch in tblPr.findall(child_tag): tblPr.remove(ch)

# ───────────── pptx shape factories ──────────────────────────
def _rect(slide, left, top, w, h, fill=None, line_color=None, line_w=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, E(left), E(top), E(w), E(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line_color:
        s.line.color.rgb = line_color
        if line_w: s.line.width = line_w
    else:
        s.line.fill.background()
    return s

def _oval(slide, left, top, size, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, E(left), E(top), E(size), E(size))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s

def _textbox(slide, left, top, w, h, text='', size=8, bold=False, color=_FALLBACK_BLACK33,
             align=PP_ALIGN.LEFT, font=_FALLBACK_FONT, wrap=True, valign='top'):
    tb = slide.shapes.add_textbox(E(left), E(top), E(w), E(h))
    tf = tb.text_frame; tf.word_wrap = wrap; tf.auto_size = None
    bp = tf._txBody.find(qn('a:bodyPr'))
    if bp is not None:
        bp.set('lIns', '0'); bp.set('tIns', '0')
        bp.set('rIns', '0'); bp.set('bIns', '0')
        if valign == 'ctr':
            bp.set('anchor', 'ctr')
    p = tf.paragraphs[0]; p.alignment = align
    p.space_before = Pt(0); p.space_after = Pt(0)
    if text:
        r = p.add_run(); r.text = text
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return tb

def _add_run(paragraph, text, size=8, bold=False, color=_FALLBACK_BLACK33, font=_FALLBACK_FONT):
    r = paragraph.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color; r.font.name = font
    return r

# ───────────── rich inline text rendering ─────────────────────
def _render_rich(paragraph, el, pt=8, skip_blocks=False):
    """Walk *el* children and emit runs with bold / colour."""
    parts: List[Tuple] = []
    if el.text:
        parts.append(('n', el.text))
    for sub in el:
        if skip_blocks and sub.tag in ('ul', 'div', 'table'):
            if sub.tail: parts.append(('n', sub.tail))
            continue
        if sub.tag in ('strong', 'b'):
            parts.append(('b', sub.text_content()))
        elif sub.tag in ('em', 'i'):
            parts.append(('i', sub.text_content()))
        elif sub.tag == 'span':
            ss = _sty(sub)
            bg = ss.get('background', '') or ss.get('background-color', '')
            if bg and ('border-radius' in str(ss) or 'display' in ss):
                cc = _parse_color(bg)
                if cc:
                    parts.append(('c', '\u25cf', cc))
                if sub.tail: parts.append(('n', sub.tail))
                continue
            txt = sub.text_content().strip()
            if txt == '\u25cf':
                c = _parse_color(ss.get('color', ''))
                if c:
                    parts.append(('c', '\u25cf', c))
                if sub.tail: parts.append(('n', sub.tail))
                continue
            if txt == '':
                if sub.tail: parts.append(('n', sub.tail))
                continue
            c = _parse_color(ss.get('color', ''))
            parts.append(('c', sub.text_content(), c))
        elif sub.tag == 'a':
            parts.append(('n', sub.text_content()))
        else:
            parts.append(('n', sub.text_content()))
        if sub.tail:
            parts.append(('n', sub.tail))

    if not parts:
        clean = el.text_content().strip()
        if clean:
            _add_run(paragraph, clean, pt)
        return

    for p in parts:
        txt = p[1]
        if skip_blocks:
            txt = txt.strip('\n')
        if not txt.strip():
            continue
        if p[0] == 'b':
            _add_run(paragraph, txt, pt, bold=True)
        elif p[0] == 'i':
            r = _add_run(paragraph, txt, pt)
            r.font.italic = True
        elif p[0] == 'c' and len(p) > 2 and p[2]:
            _add_run(paragraph, txt, pt, color=p[2])
        else:
            _add_run(paragraph, txt, pt)

# ───────────── detect circle indicator colour ─────────────────
def _circle_color(el) -> Optional[RGBColor]:
    for span in el.cssselect('span[style]'):
        ss = _sty(span)
        bg = ss.get('background', '') or ss.get('background-color', '')
        if bg:
            c = _parse_color(bg)
            if c: return c
        if ss.get('color') and '\u25cf' in (span.text or ''):
            return _parse_color(ss.get('color', ''))
    return None

# ───────────── detect progress bar structure ──────────────────
def _has_progress_bar(box_el) -> bool:
    """Check if a section-box contains progress bar divs."""
    for child in box_el:
        for inner in child:
            if not hasattr(inner, 'tag') or inner.tag != 'div': continue
            iss = _sty(inner)
            if 'border-radius' in iss and 'height' in iss and _bg_color(iss):
                return True
    return False

# ═══════════════════════════════════════════════════════════════
#  MAIN RENDERER
# ═══════════════════════════════════════════════════════════════
class SlideRenderer:
    def __init__(self, pptx_slide, html_el, ss: dict, font: str):
        self.s = pptx_slide
        self.el = html_el
        self.ss = ss
        self.font = font

    def render(self):
        self._chrome()
        self._positioned_blocks()
        self._legend()
        self._links()

    # ── chrome ─────────────────────────────────────────────────
    def _chrome(self):
        # top bar
        tbs = self.el.cssselect('.top-bar')
        if tbs:
            tb_sty = _sty(tbs[0])
            tb_h = _px(tb_sty.get('height', '')) or _px(_ss_get(self.ss, '.top-bar').get('height', '8'))
            tb_bg = _bg_color(tb_sty) or _bg_color(_ss_get(self.ss, '.top-bar')) or _FALLBACK_GREY_CC
            _rect(self.s, 0, 0, SLIDE_W_PX, tb_h, fill=tb_bg)
        # date box
        dbs = self.el.cssselect('.date-box')
        if dbs:
            db_css = _ss_get(self.ss, '.date-box')
            db_sty = _sty(dbs[0])
            db_w = _px(db_sty.get('width', '') or db_css.get('width', '100'))
            db_h = _px(db_sty.get('height', '') or db_css.get('height', '50'))
            db_top = _px(db_sty.get('top', '') or db_css.get('top', '8'))
            db_right = _px(db_sty.get('right', '') or db_css.get('right', '0'))
            db_left = SLIDE_W_PX - db_w - db_right
            db_bg = _bg_color(db_sty) or _bg_color(db_css) or _FALLBACK_GREY_CC
            db_color = _parse_color(db_sty.get('color', '') or db_css.get('color', '')) or _FALLBACK_WHITE
            db_fs = _px(db_sty.get('font-size', '') or db_css.get('font-size', '14')) * 0.75
            db_fw = db_sty.get('font-weight', '') or db_css.get('font-weight', '600')
            db_bold = db_fw not in ('400', 'normal', '')
            _rect(self.s, db_left, db_top, db_w, db_h, fill=db_bg)
            _textbox(self.s, db_left, db_top, db_w, db_h,
                     dbs[0].text_content().strip(), size=db_fs, bold=db_bold,
                     color=db_color, align=PP_ALIGN.CENTER, valign='ctr', font=self.font)
        # title
        titles = self.el.cssselect('.main-title')
        if titles:
            t = titles[0]; st = _sty(t)
            t_css = _ss_get(self.ss, '.main-title')
            fs = _px(st.get('font-size', '') or t_css.get('font-size', '42'))
            t_left = _px(st.get('left', '') or t_css.get('left', '30'))
            t_top = _px(st.get('top', '') or t_css.get('top', '20'))
            t_color = _parse_color(st.get('color', '') or t_css.get('color', '')) or _FALLBACK_TEAL
            mw = _px(st.get('max-width', '800')) if 'max-width' in st else 800
            t_fw = st.get('font-weight', '') or t_css.get('font-weight', '700')
            t_bold = t_fw not in ('400', 'normal', '')
            _textbox(self.s, t_left, t_top, mw, fs*1.4,
                     t.text_content().strip(), size=fs*0.75, bold=t_bold, color=t_color, font=self.font)
        # footer — support .footer-bar, .bottom-bar, or .footer
        fbs = self.el.cssselect('.footer-bar') or self.el.cssselect('.bottom-bar') or self.el.cssselect('.footer')
        if fbs:
            fb = fbs[0]; fb_sty = _sty(fb)
            fb_cls = 'footer-bar' if self.el.cssselect('.footer-bar') else ('bottom-bar' if self.el.cssselect('.bottom-bar') else 'footer')
            fb_css = _ss_get(self.ss, '.'+fb_cls, '.footer-bar', '.bottom-bar', '.footer')
            fb_h = _px(fb_sty.get('height', '') or fb_css.get('height', '32'))
            fb_bg = _bg_color(fb_sty) or _bg_color(fb_css) or _FALLBACK_GREY_CC
            fb_top = SLIDE_H_PX - fb_h
            _rect(self.s, 0, fb_top, SLIDE_W_PX, fb_h, fill=fb_bg)
            # page-number and logo may be children of footer OR siblings in the slide
            pn = fb.cssselect('.page-number') or self.el.cssselect('.page-number')
            if pn:
                pn_css = _ss_get(self.ss, '.page-number')
                pn_sty = _sty(pn[0])
                pn_color = _parse_color(pn_sty.get('color', '') or pn_css.get('color', '')) or _FALLBACK_WHITE
                pn_fs = _px(pn_sty.get('font-size', '') or pn_css.get('font-size', '14')) * 0.75
                _textbox(self.s, 20, fb_top, 100, fb_h,
                         pn[0].text_content().strip(), size=pn_fs, color=pn_color, valign='ctr', font=self.font)
            lg = fb.cssselect('.logo') or self.el.cssselect('.logo')
            if lg:
                lg_css = _ss_get(self.ss, '.logo')
                lg_sty = _sty(lg[0])
                lg_color = _parse_color(lg_sty.get('color', '') or lg_css.get('color', '')) or _FALLBACK_WHITE
                lg_fs = _px(lg_sty.get('font-size', '') or lg_css.get('font-size', '18')) * 0.75
                lg_fw = lg_sty.get('font-weight', '') or lg_css.get('font-weight', '700')
                lg_bold = lg_fw not in ('400', 'normal', '')
                _textbox(self.s, SLIDE_W_PX-140, fb_top, 120, fb_h,
                         lg[0].text_content().strip(), size=lg_fs, bold=lg_bold,
                         color=lg_color, align=PP_ALIGN.RIGHT, valign='ctr', font=self.font)

    # ── content dispatcher ─────────────────────────────────────
    def _positioned_blocks(self):
        # Track which elements we've already processed (e.g. sibling section-boxes)
        processed = set()

        all_children = list(self.el)
        for idx, div in enumerate(all_children):
            if not hasattr(div, 'tag') or div.tag != 'div':
                continue
            if id(div) in processed:
                continue

            st = _sty(div)
            div_cls = div.get('class', '') or ''

            # Accept position:absolute or divs with explicit top/left
            has_pos = st.get('position') == 'absolute'
            has_coords = bool(st.get('top') and st.get('left'))
            if not has_pos and not has_coords:
                continue

            # skip chrome elements handled elsewhere
            if div.cssselect('.footer-bar') or div.cssselect('.bottom-bar') or div.cssselect('.footer'): continue
            if 'footer-bar' in div_cls or 'bottom-bar' in div_cls or div_cls == 'footer': continue
            if 'page-number' in div_cls or 'logo' in div_cls:        continue
            if 'top-bar' in div_cls or 'date-box' in div_cls:        continue
            if 'main-title' in div_cls:                               continue
            if self._is_legend_div(div):                              continue
            if div.cssselect('a.link-text') and not div.cssselect('.section-header'): continue

            has_section = bool(div.cssselect('.section-header'))
            has_table   = bool(div.cssselect('table'))

            # Handle standalone section-box (sibling pattern)
            if 'section-box' in div_cls:
                # This section-box is a sibling of the header div — already handled
                # by the header div's processing below (via _find_sibling_box)
                continue

            # Handle standalone trend-box (sibling pattern)
            if 'trend-box' in div_cls:
                continue

            if has_section:
                box_els = div.cssselect('.section-box')
                box = box_els[0] if box_els else None

                # If no box as child, look for next sibling .section-box
                if box is None:
                    box = self._find_sibling_box(all_children, idx, processed)

                if box is not None and _has_progress_bar(box):
                    self._section_chrome(div, st)
                    self._planning_with_box(div, st, box)
                elif box is not None and box.cssselect('table'):
                    self._section_chrome(div, st)
                    tbl_el = box.cssselect('table')[0]
                    tbl_cls = tbl_el.get('class', '') or ''
                    tbl_cls_prefix = '.'+tbl_cls.split()[0] if tbl_cls.strip() else ''
                    td_css = _ss_get(self.ss, tbl_cls_prefix+' td') if tbl_cls_prefix else {}
                    tbl_inline = _sty(tbl_el)
                    is_dashed = any(
                        'dashed' in v for k, v in {**tbl_inline, **td_css}.items() if k.startswith('border')
                    )
                    sec_w = _resolve_width(st, default=420)
                    self._render_table(tbl_el, _px(st.get('left','0')),
                                       _px(st.get('top','0'))+20, sec_w,
                                       dashed=is_dashed)
                elif box is not None:
                    self._section_full_with_box(div, st, box)
                else:
                    # Header only, no box found — just render chrome
                    self._section_chrome(div, st)
            elif has_table:
                self._standalone_table(div, st)

    def _find_sibling_box(self, all_children, header_idx, processed) -> Optional[any]:
        """Find the next sibling .section-box or .trend-box after a section-header div."""
        for j in range(header_idx + 1, min(header_idx + 3, len(all_children))):
            sib = all_children[j]
            if not hasattr(sib, 'get'):
                continue
            sib_cls = sib.get('class', '') or ''
            if 'section-box' in sib_cls or 'trend-box' in sib_cls:
                processed.add(id(sib))
                return sib
        return None

    def _is_legend_div(self, div) -> bool:
        st = _sty(div)
        if 'bottom' not in st: return False
        if div.cssselect('.section-header'): return False
        for span in div.cssselect('span[style]'):
            ss = _sty(span)
            bg = ss.get('background', '') or ss.get('background-color', '')
            if bg and 'border-radius' in str(ss):
                return True
            if ss.get('color') and '\u25cf' in (span.text or ''):
                return True
        return False

    # ── section header + box outline ───────────────────────────
    def _section_chrome(self, div, st, ext_box=None):
        top, left, w = _px(st.get('top','0')), _px(st.get('left','0')), _resolve_width(st, default=420)
        hdr = div.cssselect('.section-header')
        hdr_sty = _sty(hdr[0]) if hdr else {}
        hdr_css = _ss_get(self.ss, '.section-header')
        sep_color = _parse_border_color(hdr_sty.get('border-top', '') or hdr_css.get('border-top', '')) or _FALLBACK_GREY_CC
        _rect(self.s, left, top, w, 1, fill=sep_color)
        titles = div.cssselect('.section-title')
        if titles:
            t_sty = _sty(titles[0])
            t_css = _ss_get(self.ss, '.section-title')
            t_color = _parse_color(t_sty.get('color', '') or t_css.get('color', '')) or _FALLBACK_TEAL
            t_fs = _px(t_sty.get('font-size', '') or t_css.get('font-size', '13')) * 0.75
            t_fw = t_sty.get('font-weight', '') or t_css.get('font-weight', '700')
            t_bold = t_fw not in ('400', 'normal', '')
            _textbox(self.s, left, top+2, w, 16,
                     titles[0].text_content().strip(), size=t_fs, bold=t_bold, color=t_color, font=self.font)
        # Box can be a child or passed externally (sibling pattern)
        box_el = ext_box
        if box_el is None:
            box_els = div.cssselect('.section-box')
            box_el = box_els[0] if box_els else None
        if box_el is not None:
            box_sty = _sty(box_el)
            box_css = _ss_get(self.ss, '.section-box')
            bh = _px(box_sty.get('height', '') or box_css.get('height', '80'))
            box_bg = _bg_color(box_sty) or _bg_color(box_css) or _FALLBACK_WHITE
            box_border = _parse_border_color(box_sty.get('border', '') or box_css.get('border', '')) or _FALLBACK_GREY_CC
            # Use box's own top/left if it has them (sibling pattern), else offset from header
            box_top_px = _px(box_sty.get('top', ''))
            box_left_px = _px(box_sty.get('left', ''))
            if box_top_px > 0:
                _rect(self.s, box_left_px or left, box_top_px, w, bh, fill=box_bg, line_color=box_border, line_w=Pt(0.75))
            else:
                _rect(self.s, left, top+20, w, bh, fill=box_bg, line_color=box_border, line_w=Pt(0.75))

    # ── full section (header + box + content) ──────────────────
    def _section_full(self, div, st):
        self._section_full_with_box(div, st, None)

    def _section_full_with_box(self, div, st, ext_box):
        self._section_chrome(div, st, ext_box)
        top, left, w = _px(st.get('top','0')), _px(st.get('left','0')), _resolve_width(st, default=420)
        if ext_box is not None:
            box = ext_box
        else:
            box_els = div.cssselect('.section-box')
            if not box_els: return
            box = box_els[0]
        # Use box's own top if available (sibling pattern)
        box_sty = _sty(box)
        box_top_px = _px(box_sty.get('top', ''))
        box_top = box_top_px if box_top_px > 0 else top + 20

        tables = box.cssselect('table')
        if tables:
            tbl_cls = tables[0].get('class', '') or ''
            tbl_cls_prefix = '.'+tbl_cls.split()[0] if tbl_cls.strip() else ''
            td_css = _ss_get(self.ss, tbl_cls_prefix+' td') if tbl_cls_prefix else {}
            tbl_inline = _sty(tables[0])
            is_dashed = any(
                'dashed' in v for k, v in {**tbl_inline, **td_css}.items() if k.startswith('border')
            )
            self._render_table(tables[0], left+2, box_top+2, w-4, dashed=is_dashed)
            return

        # Check if box IS a trend-box or contains one
        box_cls = box.get('class', '') or ''
        trend_el = box if 'trend-box' in box_cls else None
        if trend_el is None:
            trend_els = box.cssselect('.trend-box')
            trend_el = trend_els[0] if trend_els else None
        if trend_el is not None:
            ti_css = _ss_get(self.ss, '.trend-item')
            ti_fs = _px(ti_css.get('font-size', '14')) * 0.75
            ti_fw = ti_css.get('font-weight', '600')
            ti_bold = ti_fw not in ('400', 'normal', '')
            ti_color = _parse_color(ti_css.get('color', '')) or _FALLBACK_BLACK33
            ti_gap = _px(_ss_get(self.ss, '.trend-box').get('gap', '30'))
            # Use trend-box's own position if available
            trend_sty = _sty(trend_el)
            trend_top = _px(trend_sty.get('top', ''))
            trend_left = _px(trend_sty.get('left', ''))
            render_y = trend_top if trend_top > 0 else box_top + 10
            render_x = (trend_left or left) + 8
            for item in trend_el.cssselect('.trend-item'):
                item_sty = _sty(item)
                item_color = _parse_color(item_sty.get('color', '')) or ti_color
                item_fs = _px(item_sty.get('font-size', '')) * 0.75 if item_sty.get('font-size') else ti_fs
                _textbox(self.s, render_x, render_y, 80, 20,
                         item.text_content().strip(), size=item_fs, bold=ti_bold,
                         color=item_color, font=self.font)
                render_x += 80 + ti_gap
            return

        y = box_top + 6
        y = self._render_box_content(box, left, y, w)

    # ── recursive box content renderer ────────────────────────
    _INLINE_TAGS = {'strong', 'b', 'em', 'i', 'span', 'a', 'br', 'sub', 'sup'}

    def _render_box_content(self, parent, left, y, w, indent=0, fs_pt=8):
        x_base = left + 8 + indent
        w_inner = w - 16 - indent
        LINE_H = 13

        bi_before = _ss_get(self.ss, '.bullet-item::before')
        bullet_color = _parse_color(bi_before.get('color', '')) or _FALLBACK_RED_BULL
        bullet_char = bi_before.get('content', '').strip('"\'').replace('\\25aa', '\u25aa').replace('\\2022', '\u2022') or '\u25aa'
        if bullet_char.startswith('\\') and len(bullet_char) <= 5:
            try: bullet_char = chr(int(bullet_char[1:], 16))
            except: pass
        bullet_fs = _px(bi_before.get('font-size', '10')) * 0.75

        for child in parent:
            if child.tag in self._INLINE_TAGS:
                continue
            tag = child.tag
            cls = child.get('class', '') or ''
            cst = _sty(child)

            mt = _px(cst.get('margin-top', '0'))
            mb = _px(cst.get('margin-bottom', '0'))
            y += mt

            lh_str = cst.get('line-height', '')
            lh = LINE_H
            if lh_str:
                lh_val = _px(lh_str)
                if lh_val > 0 and lh_val < 5:
                    lh = int(LINE_H * lh_val)
                elif lh_val >= 5:
                    lh = int(lh_val)

            if 'budget-label' in cls or 'sub-label' in cls:
                bl_css = _ss_get(self.ss, '.budget-label', '.sub-label')
                bl_fs = _px(cst.get('font-size', '') or bl_css.get('font-size', '12')) * 0.75
                bl_fw = cst.get('font-weight', '') or bl_css.get('font-weight', '700')
                bl_bold = bl_fw not in ('400', 'normal', '')
                bl_color = _parse_color(cst.get('color', '') or bl_css.get('color', '')) or _FALLBACK_BLACK33
                _textbox(self.s, x_base, y, w_inner, 14,
                         child.text_content().strip(), size=bl_fs, bold=bl_bold, color=bl_color, font=self.font)
                y += 14 + mb; continue

            if 'bullet-item' in cls:
                bi_css = _ss_get(self.ss, '.bullet-item')
                fs_css = _px(cst.get('font-size', '') or bi_css.get('font-size', '11'))
                fpt = fs_css * 0.75
                item_mb = _px(cst.get('margin-bottom', '') or bi_css.get('margin-bottom', '4'))
                _textbox(self.s, x_base, y, 10, 12, bullet_char, size=bullet_fs, color=bullet_color, font=self.font)
                tb = _textbox(self.s, x_base+12, y, w_inner-12, 14, font=self.font)
                _render_rich(tb.text_frame.paragraphs[0], child, fpt)
                y += max(14, int(fpt*1.8)) + item_mb; continue

            if tag in ('ul', 'ol'):
                ul_margin = _px(cst.get('margin-left', '0'))
                li_x = x_base + ul_margin
                li_w = w_inner - ul_margin
                li_num = 0
                for li in child:
                    if li.tag != 'li': continue
                    li_num += 1
                    txt = li.text_content().strip()
                    if not txt: continue
                    li_sty = _sty(li)
                    li_color = _parse_color(li_sty.get('color', '')) or _FALLBACK_BLACK33
                    if tag == 'ol':
                        marker = f'{li_num}.'
                        _textbox(self.s, li_x, y, 18, lh, marker, size=8, color=li_color, font=self.font)
                        tb = _textbox(self.s, li_x+18, y, li_w-18, lh, font=self.font)
                    else:
                        _textbox(self.s, li_x, y, 10, lh, '\u2022', size=8, color=li_color, font=self.font)
                        tb = _textbox(self.s, li_x+12, y, li_w-12, lh, font=self.font)
                    _render_rich(tb.text_frame.paragraphs[0], li, fs_pt)
                    y += lh
                    # Handle nested lists inside <li>
                    for nested in li:
                        if nested.tag in ('ul', 'ol'):
                            y = self._render_box_content(nested, left, y, w, indent + 16, fs_pt)
                y += mb
                continue

            if tag == 'p':
                txt = child.text_content().strip()
                if txt:
                    tb = _textbox(self.s, x_base, y, w_inner, 14, font=self.font)
                    _render_rich(tb.text_frame.paragraphs[0], child, fs_pt)
                    y += 14
                y += mb
                continue

            if tag == 'div':
                has_blocks = any(c.tag in ('p', 'ul', 'div', 'table') for c in child)
                if has_blocks:
                    child_fs = fs_pt
                    fs_str = cst.get('font-size', '')
                    if fs_str:
                        child_fs = _px(fs_str) * 0.75
                    child_indent = indent + _px(cst.get('margin-left', '0'))
                    y = self._render_box_content(child, left, y, w, child_indent, child_fs)
                else:
                    txt = child.text_content().strip()
                    if txt:
                        tb = _textbox(self.s, x_base, y, w_inner, 14, font=self.font)
                        _render_rich(tb.text_frame.paragraphs[0], child, fs_pt)
                        y += 14
                y += mb
                continue

        return y

    # ── planning section ──────────────────────────────────────
    def _planning(self, div, st):
        self._planning_with_box(div, st, None)

    def _planning_with_box(self, div, st, ext_box):
        top, left, w = _px(st.get('top','0')), _px(st.get('left','0')), _resolve_width(st, default=900)
        if ext_box is not None:
            box = ext_box
        else:
            box_els = div.cssselect('.section-box')
            if not box_els: return
            box = box_els[0]
        box_sty = _sty(box)
        box_top_px = _px(box_sty.get('top', ''))
        y = (box_top_px if box_top_px > 0 else top + 28)

        for child in box:
            drew_bar = False
            for inner in child:
                if not hasattr(inner, 'tag') or inner.tag != 'div': continue
                iss = _sty(inner)
                if 'border-radius' not in iss or 'height' not in iss: continue
                inner_bg = _bg_color(iss)
                if not inner_bg: continue
                fill_children = [fd for fd in inner if fd.tag == 'div' and _bg_color(_sty(fd))]
                has_span = any(sp.tag == 'span' for sp in inner)
                if not fill_children and not has_span: continue
                drew_bar = True
                bar_h = _px(iss.get('height', '16'))
                bar_w = w - 24
                bg = self.s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                        E(left+12), E(y), E(bar_w), E(bar_h))
                bg.fill.solid(); bg.fill.fore_color.rgb = inner_bg; bg.line.fill.background()
                for fd in fill_children:
                    fds = _sty(fd)
                    fill_color = _bg_color(fds)
                    if not fill_color: continue
                    pct = _pct(fds.get('width','0'))
                    fw = bar_w * pct / 100
                    if fw > 0:
                        fb = self.s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                E(left+12), E(y), E(fw), E(bar_h))
                        fb.fill.solid(); fb.fill.fore_color.rgb = fill_color; fb.line.fill.background()
                for sp in inner:
                    if sp.tag != 'span': continue
                    stxt = sp.text_content().strip()
                    if '%' in stxt:
                        sp_sty = _sty(sp)
                        sp_fs = _px(sp_sty.get('font-size', '11')) * 0.75 if sp_sty.get('font-size') else 8
                        sp_color = _parse_color(sp_sty.get('color', '')) or _FALLBACK_BLACK33
                        _textbox(self.s, left+bar_w-60, y, 72, bar_h,
                                 stxt, size=sp_fs, color=sp_color, align=PP_ALIGN.RIGHT, font=self.font)
                y += bar_h + 8
            if drew_bar: continue

            txt = child.text_content().strip()
            if not txt or child.tag != 'div': continue

            has_blocks = any(c.tag in ('p', 'ul', 'div', 'table') for c in child)
            if has_blocks:
                tb = _textbox(self.s, left+12, y, w-24, 14, font=self.font)
                _render_rich(tb.text_frame.paragraphs[0], child, 8, skip_blocks=True)
                y += 16
                y = self._render_box_content(child, left, y, w, indent=12, fs_pt=8)
            else:
                tb = _textbox(self.s, left+12, y, w-24, 14, font=self.font)
                _render_rich(tb.text_frame.paragraphs[0], child, 8)
                y += 16

    # ── standalone table (summary slide) ───────────────────────
    def _standalone_table(self, div, st):
        top, left = _px(st.get('top','0')), _px(st.get('left','0'))
        w = _resolve_width(st, default=0)
        if w == 0:
            w = SLIDE_W_PX - left - 20  # fallback: fill remaining space
        tables = div.cssselect('table')
        if tables:
            # Also check if table itself has a width
            tbl_w = _resolve_width(_sty(tables[0]), default=0, container_w=w)
            if tbl_w > 0:
                w = tbl_w
            self._render_table(tables[0], left, top, w)

    # ── generic HTML table → pptx table ───────────────────────
    def _render_table(self, table_el, left, top, width, dashed=False):
        trs = list(table_el.cssselect('tr'))
        if not trs: return
        # Find max columns across all rows (handles colspan/irregular rows)
        n_cols = max((len(tr.cssselect('th, td')) for tr in trs), default=0)
        n_rows = len(trs)
        if not n_cols or not n_rows: return
        if width <= 0: width = SLIDE_W_PX - left - 20

        tbl_sty = _sty(table_el)
        tbl_fs_px = _px(tbl_sty.get('font-size', '11'))
        tbl_border_color = _parse_border_color(tbl_sty.get('border', '')) or _FALLBACK_GREY_CC

        # Also check table-level width
        tbl_w = _resolve_width(tbl_sty, default=0, container_w=width)
        if tbl_w > 0:
            width = tbl_w

        first_cells = trs[0].cssselect('th, td')
        explicit_w = []
        for c in first_cells:
            cs = _sty(c)
            # Try inline width (px or %)
            cw = _resolve_width(cs, default=0, container_w=width)
            # Check stylesheet for column widths
            if cw == 0:
                cls = c.get('class', '') or ''
                if cls:
                    for cn in cls.split():
                        cls_css = _ss_get(self.ss, '.'+cn)
                        cw = _resolve_width(cls_css, default=0, container_w=width)
                        if cw > 0:
                            break
            explicit_w.append(cw)
        total_explicit = sum(explicit_w)
        remaining = width - total_explicit
        n_auto = sum(1 for w in explicit_w if w == 0)
        auto_w = remaining / max(n_auto, 1) if n_auto else 0
        col_widths = [w if w > 0 else auto_w for w in explicit_w]

        tbl_class = table_el.get('class', '') or ''
        tbl_cls_prefix = '.'+tbl_class.split()[0] if tbl_class.strip() else ''

        th_css = _ss_get(self.ss, tbl_cls_prefix+' th') if tbl_cls_prefix else {}
        td_css = _ss_get(self.ss, tbl_cls_prefix+' td') if tbl_cls_prefix else {}
        # detect dashed from any border property
        def _has_dashed(css):
            for k, v in css.items():
                if k.startswith('border') and 'dashed' in v:
                    return True
            return False
        if not dashed and td_css:
            dashed = _has_dashed(td_css)
        if not dashed and th_css:
            dashed = _has_dashed(th_css)

        # resolve border color from any border property
        def _any_border_color(css):
            for k in ('border', 'border-top', 'border-bottom', 'border-left', 'border-right'):
                c = _parse_border_color(css.get(k, ''))
                if c: return c
            return None
        cell_border_css = _any_border_color(td_css) or _any_border_color(th_css) or tbl_border_color

        row_h = 22
        shape = self.s.shapes.add_table(n_rows, n_cols,
                    E(left), E(top), E(width), E(n_rows * row_h + 4))
        tbl = shape.table
        _nuke_table_theme(shape)

        for ci, cw in enumerate(col_widths):
            if ci < n_cols:
                tbl.columns[ci].width = E(cw)

        for ri, tr in enumerate(trs):
            cells = tr.cssselect('th, td')
            tr_sty = _sty(tr)
            tr_bg = _parse_color(tr_sty.get('background', ''))
            tr_color = _parse_color(tr_sty.get('color', ''))
            for ci, td in enumerate(cells):
                if ci >= n_cols: break
                cell = tbl.cell(ri, ci)
                ds = _sty(td)
                cls = td.get('class', '') or ''

                cls_css = {}
                if cls:
                    for c in cls.split():
                        cls_css.update(_ss_get(self.ss, '.'+c))

                txt = td.text_content().strip()
                cc = _circle_color(td)
                if cc: txt = '\u25cf'

                cell_fs = _px(ds.get('font-size', '')) or _px(cls_css.get('font-size', '')) or \
                          _px((th_css if td.tag == 'th' else td_css).get('font-size', '')) or tbl_fs_px
                fs_pt = cell_fs * 0.75
                if fs_pt < 6: fs_pt = 8

                align = ds.get('text-align', '') or cls_css.get('text-align', '') or \
                        (th_css if td.tag == 'th' else td_css).get('text-align', '') or 'left'

                fw = ds.get('font-weight', '') or cls_css.get('font-weight', '') or \
                     (th_css if td.tag == 'th' else td_css).get('font-weight', '')

                tf = cell.text_frame; tf.clear()
                p = tf.paragraphs[0]
                p.space_before = Pt(0); p.space_after = Pt(0)
                p.alignment = {'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)

                r = p.add_run(); r.font.name = self.font
                if cc:
                    r.text = '\u25cf'; r.font.size = Pt(10); r.font.color.rgb = cc
                else:
                    r.text = txt; r.font.size = Pt(fs_pt)
                    cell_color = _parse_color(ds.get('color', '')) or _parse_color(cls_css.get('color', '')) or tr_color or _FALLBACK_BLACK33
                    r.font.color.rgb = cell_color

                if fw:
                    r.font.bold = fw not in ('400', 'normal')
                elif td.tag == 'th':
                    r.font.bold = True

                cell_bg = _bg_color(ds) or _bg_color(cls_css)
                if td.tag == 'th':
                    bg = cell_bg or tr_bg
                    if bg:
                        _cell_fill(cell, bg)
                        if not _is_light(bg):
                            r.font.color.rgb = _FALLBACK_WHITE
                elif tr_bg:
                    _cell_fill(cell, tr_bg)
                elif cell_bg:
                    _cell_fill(cell, cell_bg)

                _cell_border(cell, cell_border_css, 6350, 'dashed' if dashed else 'solid')
                cell.margin_left = E(4); cell.margin_right = E(4)
                cell.margin_top  = E(2); cell.margin_bottom = E(2)

    # ── legend (summary slide) ─────────────────────────────────
    def _legend(self):
        for el in self.el.cssselect('div[style]'):
            if not self._is_legend_div(el): continue
            st = _sty(el)
            bottom = _px(st.get('bottom','50'))
            lx = _px(st.get('left','30'))
            leg_fs = _px(st.get('font-size', '11')) * 0.75
            leg_color = _parse_color(st.get('color', '')) or _FALLBACK_GREY66
            ty = SLIDE_H_PX - bottom - 20

            tb = _textbox(self.s, lx, ty, 800, 20, font=self.font)
            _render_rich(tb.text_frame.paragraphs[0], el, leg_fs)
            for run in tb.text_frame.paragraphs[0].runs:
                if run.font.color.rgb in (None, _FALLBACK_BLACK33):
                    run.font.color.rgb = leg_color

    # ── links ──────────────────────────────────────────────────
    def _links(self):
        link_css = _ss_get(self.ss, '.link-text')
        link_color = _parse_color(link_css.get('color', '')) or _FALLBACK_TEAL
        link_fs = _px(link_css.get('font-size', '12')) * 0.75

        # Links inside positioned divs
        for el in self.el.cssselect('div[style]'):
            st = _sty(el)
            if st.get('position') != 'absolute': continue
            links = el.cssselect('a.link-text')
            if not links or el.cssselect('.section-header'): continue
            bottom = _px(st.get('bottom','60'))
            lx = _px(st.get('left','30'))
            ty = SLIDE_H_PX - bottom - 15
            for a in links:
                a_sty = _sty(a)
                a_color = _parse_color(a_sty.get('color', '')) or link_color
                a_fs = _px(a_sty.get('font-size', '')) * 0.75 if a_sty.get('font-size') else link_fs
                tb = _textbox(self.s, lx, ty, 300, 15,
                              a.text_content().strip(), size=a_fs, color=a_color, font=self.font)
                tb.text_frame.paragraphs[0].runs[0].font.underline = True

        # Links as direct children of the slide (not wrapped in a div)
        for a in self.el.cssselect('a.link-text'):
            parent = a.getparent()
            if parent is not None and parent is not self.el:
                continue  # already handled above (inside a div)
            a_sty = _sty(a)
            a_top = _px(a_sty.get('top', ''))
            a_left = _px(a_sty.get('left', '30'))
            a_bottom = _px(a_sty.get('bottom', ''))
            if a_top > 0:
                ty = a_top
            elif a_bottom > 0:
                ty = SLIDE_H_PX - a_bottom - 15
            else:
                ty = SLIDE_H_PX - 75
            a_color = _parse_color(a_sty.get('color', '')) or link_color
            a_fs = _px(a_sty.get('font-size', '')) * 0.75 if a_sty.get('font-size') else link_fs
            tb = _textbox(self.s, a_left, ty, 300, 15,
                          a.text_content().strip(), size=a_fs, color=a_color, font=self.font)
            tb.text_frame.paragraphs[0].runs[0].font.underline = True


# ═══════════════════════════════════════════════════════════════
#  PUBLIC API
# ═══════════════════════════════════════════════════════════════
def _resolve_font(ss: dict) -> str:
    """Extract font-family from body or .slide stylesheet rules."""
    for sel in ('body', '.slide'):
        ff = ss.get(sel, {}).get('font-family', '')
        if ff:
            first = ff.split(',')[0].strip().strip("'\"")
            if first and first != '-apple-system':
                return first
            parts = [p.strip().strip("'\"") for p in ff.split(',')]
            for p in parts:
                if p and p not in ('-apple-system', 'BlinkMacSystemFont'):
                    return p
    return _FALLBACK_FONT


def html_to_pptx(html_content: str) -> bytes:
    """
    Convert populated HTML slides to PPTX bytes.

    Args:
        html_content: Full HTML string with div.slide elements

    Returns:
        PPTX file content as bytes
    """
    doc = lxml_html.fromstring(html_content)

    ss = _parse_stylesheet(doc)
    font = _resolve_font(ss)

    slide_els = doc.cssselect('div.slide')
    if not slide_els:
        raise ValueError("No slides found in HTML (expected div.slide elements)")

    prs = Presentation()
    prs.slide_width  = Emu(int(SLIDE_W_IN * 914400))
    prs.slide_height = Emu(int(SLIDE_H_IN * 914400))
    blank = prs.slide_layouts[6]

    for i, el in enumerate(slide_els):
        sl = prs.slides.add_slide(blank)
        try:
            SlideRenderer(sl, el, ss, font).render()
            print(f'  [PPTX] [{i+1}/{len(slide_els)}] rendered')
        except Exception as exc:
            print(f'  [PPTX] [{i+1}/{len(slide_els)}] ERROR: {exc}')
            # Add a textbox with error message so the slide isn't blank
            _textbox(sl, 30, 250, 900, 40,
                     f'Slide rendering error: {exc}', size=10, color=_FALLBACK_RED_BULL)

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)

    print(f'[PPTX] Generated {len(slide_els)} slides')
    return buffer.read()
